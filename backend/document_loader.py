import os
import fitz
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from langchain_text_splitters import RecursiveCharacterTextSplitter

DOCUMENTS_FOLDER = "documents"

splitter = RecursiveCharacterTextSplitter(
    chunk_size=800,
    chunk_overlap=150
)


def clean_text(text):
    return " ".join(text.split())


def load_documents():

    chunks = []
    metadata = []

    if not os.path.exists(DOCUMENTS_FOLDER):
        os.makedirs(DOCUMENTS_FOLDER)

    for file in os.listdir(DOCUMENTS_FOLDER):

        file_path = os.path.join(DOCUMENTS_FOLDER, file)

        # =====================================================
        # PDF
        # =====================================================

        if file.lower().endswith(".pdf"):

            pdf = fitz.open(file_path)

            total_pages = len(pdf)

            for page_number, page in enumerate(pdf):

                text = clean_text(page.get_text())

                if not text:
                    continue

                page_chunks = splitter.split_text(text)

                for chunk in page_chunks:

                    chunks.append(chunk)

                    metadata.append({

                        "document": file,

                        "file_type": "PDF",

                        "reference": f"Page {page_number + 1} of {total_pages}",

                        "page_number": page_number + 1,

                        "heading": "PDF Content",

                        "chunk": chunk

                    })

            pdf.close()

        # =====================================================
        # WORD
        # =====================================================

        elif file.lower().endswith(".docx"):

            document = Document(file_path)

            current_heading = "Document"

            for para in document.paragraphs:

                text = clean_text(para.text)

                if not text:
                    continue

                if para.style.name.startswith("Heading"):

                    current_heading = text

                page_chunks = splitter.split_text(text)

                for chunk in page_chunks:

                    chunks.append(chunk)

                    metadata.append({

                        "document": file,

                        "file_type": "Word",

                        "reference": "Section",

                        "page_number": None,

                        "heading": current_heading,

                        "chunk": chunk

                    })

        # =====================================================
        # POWERPOINT
        # =====================================================

        elif file.lower().endswith(".pptx"):

            presentation = Presentation(file_path)

            total_slides = len(presentation.slides)

            for slide_number, slide in enumerate(presentation.slides):

                slide_title = "Untitled Slide"

                slide_text = ""

                if slide.shapes.title:

                    title = clean_text(slide.shapes.title.text)

                    if title:

                        slide_title = title

                for shape in slide.shapes:

                    if hasattr(shape, "text"):

                        slide_text += shape.text + "\n"

                slide_text = clean_text(slide_text)

                if not slide_text:
                    continue

                page_chunks = splitter.split_text(slide_text)

                for chunk in page_chunks:

                    chunks.append(chunk)

                    metadata.append({

                        "document": file,

                        "file_type": "PowerPoint",

                        "reference": f"Slide {slide_number + 1} of {total_slides}",

                        "page_number": None,

                        "heading": slide_title,

                        "chunk": chunk

                    })

        # =====================================================
        # EXCEL
        # =====================================================

        elif file.lower().endswith(".xlsx"):

            workbook = load_workbook(file_path)

            total_sheets = len(workbook.sheetnames)

            for sheet_number, worksheet in enumerate(workbook.worksheets):

                sheet_name = worksheet.title

                rows = []

                for row in worksheet.iter_rows(values_only=True):

                    values = [

                        str(cell)

                        for cell in row

                        if cell is not None

                    ]

                    if values:

                        rows.append(" ".join(values))

                sheet_text = clean_text("\n".join(rows))

                if not sheet_text:
                    continue

                page_chunks = splitter.split_text(sheet_text)

                for chunk in page_chunks:

                    chunks.append(chunk)

                    metadata.append({

                        "document": file,

                        "file_type": "Excel",

                        "reference": f"Sheet {sheet_number + 1} of {total_sheets}",

                        "page_number": None,

                        "heading": sheet_name,

                        "chunk": chunk

                    })

        # =====================================================
        # TEXT
        # =====================================================

        elif file.lower().endswith(".txt"):

            with open(

                file_path,

                "r",

                encoding="utf-8"

            ) as f:

                text = clean_text(f.read())

            if not text:
                continue

            page_chunks = splitter.split_text(text)

            for chunk in page_chunks:

                chunks.append(chunk)

                metadata.append({

                    "document": file,

                    "file_type": "Text",

                    "reference": "Text File",

                    "page_number": None,

                    "heading": "General",

                    "chunk": chunk

                })

    return chunks, metadata


def load_single_document(file_path):

    chunks = []
    metadata = []

    filename = os.path.basename(file_path)
    extension = filename.lower().split(".")[-1]

    print(f"Loading uploaded document: {filename}")

    single_splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=100
    )

    try:

        if extension == "pdf":

            doc = fitz.open(file_path)

            total_pages = len(doc)

            for page_number, page in enumerate(doc):

                text = clean_text(page.get_text())

                if not text:
                    continue

                document_chunks = single_splitter.split_text(text)

                for chunk in document_chunks:

                    chunks.append(chunk)

                    metadata.append({

                        "document": filename,

                        "file_type": "PDF",

                        "reference": f"Page {page_number + 1} of {total_pages}",

                        "page_number": page_number + 1,

                        "heading": "PDF Content",

                        "chunk": chunk

                    })

            doc.close()

        elif extension == "docx":

            document = Document(file_path)

            current_heading = "Document"

            for para in document.paragraphs:

                text = clean_text(para.text)

                if not text:
                    continue

                if para.style.name.startswith("Heading"):
                    current_heading = text

                document_chunks = single_splitter.split_text(text)

                for chunk in document_chunks:

                    chunks.append(chunk)

                    metadata.append({

                        "document": filename,

                        "file_type": "Word",

                        "reference": "Section",

                        "page_number": None,

                        "heading": current_heading,

                        "chunk": chunk

                    })

        elif extension == "pptx":

            presentation = Presentation(file_path)

            total_slides = len(presentation.slides)

            for slide_number, slide in enumerate(presentation.slides):

                slide_title = "Untitled Slide"

                slide_text = ""

                if slide.shapes.title:

                    title = clean_text(slide.shapes.title.text)

                    if title:
                        slide_title = title

                for shape in slide.shapes:

                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"

                slide_text = clean_text(slide_text)

                if not slide_text:
                    continue

                document_chunks = single_splitter.split_text(slide_text)

                for chunk in document_chunks:

                    chunks.append(chunk)

                    metadata.append({

                        "document": filename,

                        "file_type": "PowerPoint",

                        "reference": f"Slide {slide_number + 1} of {total_slides}",

                        "page_number": None,

                        "heading": slide_title,

                        "chunk": chunk

                    })

        elif extension == "xlsx":

            workbook = load_workbook(file_path)

            total_sheets = len(workbook.sheetnames)

            for sheet_number, worksheet in enumerate(workbook.worksheets):

                sheet_name = worksheet.title

                rows = []

                for row in worksheet.iter_rows(values_only=True):

                    values = [
                        str(cell)
                        for cell in row
                        if cell is not None
                    ]

                    if values:
                        rows.append(" ".join(values))

                sheet_text = clean_text("\n".join(rows))

                if not sheet_text:
                    continue

                document_chunks = single_splitter.split_text(sheet_text)

                for chunk in document_chunks:

                    chunks.append(chunk)

                    metadata.append({

                        "document": filename,

                        "file_type": "Excel",

                        "reference": f"Sheet {sheet_number + 1} of {total_sheets}",

                        "page_number": None,

                        "heading": sheet_name,

                        "chunk": chunk

                    })

        elif extension == "txt":

            with open(file_path, "r", encoding="utf-8") as file:
                text = clean_text(file.read())

            document_chunks = single_splitter.split_text(text)

            for chunk in document_chunks:

                chunks.append(chunk)

                metadata.append({

                    "document": filename,

                    "file_type": "Text",

                    "reference": "Text File",

                    "page_number": None,

                    "heading": "General",

                    "chunk": chunk

                })

        else:

            print(f"Unsupported file: {filename}")
            return [], []

    except Exception as e:

        print(f"Error reading {filename}: {e}")
        return [], []

    print(f"{len(chunks)} chunks created.")

    return chunks, metadata